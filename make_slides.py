"""Generate modern narrative-focused capstone presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Palette ──────────────────────────────────────────────────────────────────
BG         = RGBColor(0x0F, 0x0F, 0x0F)   # near-black
PANEL      = RGBColor(0x1A, 0x1A, 0x1A)   # dark card
ACCENT     = RGBColor(0x6C, 0x63, 0xFF)   # purple
ACCENT2    = RGBColor(0x00, 0xC9, 0xA7)   # teal  (positive/upgrade)
WARN       = RGBColor(0xFF, 0x6B, 0x6B)   # red   (problem/limitation)
GOLD       = RGBColor(0xFF, 0xD1, 0x66)   # gold  (key line highlight)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE   = RGBColor(0xE8, 0xE8, 0xE8)
MUTED      = RGBColor(0x88, 0x88, 0x99)
# ─────────────────────────────────────────────────────────────────────────────


def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, radius=False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font="Calibri", wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.font.italic = italic
    return tf


def bullets(slide, l, t, w, h, items, size=17, color=OFFWHITE,
            dot_color=ACCENT, spacing=10):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)
    return tf


def key_line(slide, text, y, color=GOLD):
    """Big bold key statement — the one thing to remember from the slide."""
    rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.75), PANEL)
    txt(slide, Inches(1.1), y + Inches(0.08), Inches(11.1), Inches(0.6),
        text, size=22, color=color, bold=True, align=PP_ALIGN.CENTER, font="Calibri")


def stage_badge(slide, label, color=ACCENT):
    rect(slide, Inches(0.8), Inches(0.72), Inches(1.6), Inches(0.36), color)
    txt(slide, Inches(0.82), Inches(0.74), Inches(1.56), Inches(0.32),
        label, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def divider(slide, y, color=ACCENT, alpha_sim=False):
    rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.04), color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

# Left accent bar
rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT)

# Subtle grid lines (decorative)
for i in range(1, 5):
    rect(s, Inches(0.06), Inches(i * 1.5), Inches(13.333), Inches(0.01),
         RGBColor(0x22, 0x22, 0x22))

txt(s, Inches(1.0), Inches(1.8), Inches(11), Inches(0.5),
    "CAPSTONE PROJECT", size=14, color=ACCENT, bold=True, font="Calibri")

txt(s, Inches(1.0), Inches(2.3), Inches(11), Inches(1.6),
    "Personal Wardrobe Intelligence", size=46, bold=True,
    align=PP_ALIGN.LEFT, font="Calibri")

rect(s, Inches(1.0), Inches(4.0), Inches(5), Inches(0.04), ACCENT2)

txt(s, Inches(1.0), Inches(4.3), Inches(10), Inches(0.6),
    "Visual Similarity Retrieval  ·  Preference Reranking  ·  FashionCLIP + FAISS",
    size=19, color=MUTED, font="Calibri")

txt(s, Inches(1.0), Inches(5.8), Inches(6), Inches(0.4),
    "April 2026", size=16, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5), WARN)

txt(s, Inches(0.9), Inches(0.55), Inches(10), Inches(0.7),
    "The Problem", size=36, bold=True)
divider(s, Inches(1.35), WARN)

txt(s, Inches(0.9), Inches(1.65), Inches(11), Inches(0.6),
    "Finding similar clothing is slow, subjective, and tool-less.", size=22, color=OFFWHITE)

bullets(s, Inches(0.9), Inches(2.5), Inches(11), Inches(3.5), [
    "Browsing a wardrobe manually relies on memory and guesswork",
    "Keyword search fails for visual style — \"navy blue structured blazer\" returns nothing useful",
    "No existing free tool answers: \"what do I own that looks like this?\"",
    "And no tool lets you say: \"similar to this, but more casual, no logos\"",
], size=20, dot_color=WARN, color=OFFWHITE, spacing=14)

key_line(s, "\"Visual style can't be typed — it needs to be shown.\"", Inches(6.35), WARN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — STAGE 1: BASELINE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5), MUTED)
stage_badge(s, "STAGE 1", MUTED)

txt(s, Inches(0.9), Inches(0.55), Inches(10), Inches(0.7),
    "Baseline: CLIP Retrieval", size=36, bold=True)
divider(s, Inches(1.35), MUTED)

txt(s, Inches(0.9), Inches(1.6), Inches(11), Inches(0.55),
    "\"I started with a basic CLIP-based retrieval system.\"",
    size=20, color=MUTED, italic=True)

# Two columns: what it did / what it couldn't do
rect(s, Inches(0.9), Inches(2.4), Inches(5.4), Inches(3.0), PANEL)
txt(s, Inches(1.1), Inches(2.55), Inches(5.0), Inches(0.45),
    "What it did", size=17, color=ACCENT2, bold=True)
bullets(s, Inches(1.1), Inches(3.1), Inches(5.0), Inches(2.2), [
    "Upload image → get 5 similar items",
    "CLIP encodes images into 512-dim embeddings",
    "Cosine similarity over ~400 catalog images",
    "Live Streamlit web demo",
], size=16, color=OFFWHITE, spacing=10)

rect(s, Inches(7.0), Inches(2.4), Inches(5.4), Inches(3.0), PANEL)
txt(s, Inches(7.2), Inches(2.55), Inches(5.0), Inches(0.45),
    "Where it fell short", size=17, color=WARN, bold=True)
bullets(s, Inches(7.2), Inches(3.1), Inches(5.0), Inches(2.2), [
    "CLIP trained on generic internet data",
    "Struggled with specific garments",
    "Puffer jackets, cardigans → poor results",
    "No way to guide or filter results",
], size=16, color=OFFWHITE, spacing=10)

key_line(s, "\"It worked — but it wasn't fashion-aware.\"", Inches(6.35))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — STAGE 2: PREFERENCE RERANKING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT)
stage_badge(s, "STAGE 2", ACCENT)

txt(s, Inches(0.9), Inches(0.55), Inches(10), Inches(0.7),
    "Adding User Preferences", size=36, bold=True)
divider(s, Inches(1.35), ACCENT)

txt(s, Inches(0.9), Inches(1.6), Inches(11), Inches(0.55),
    "\"I added user preferences to make results controllable.\"",
    size=20, color=MUTED, italic=True)

# Formula box
rect(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.8), PANEL)
txt(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(0.6),
    "final_score  =  base_score  +  0.15 × goal_bonus  −  0.15 × avoid_penalty",
    size=20, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER, font="Consolas")

bullets(s, Inches(0.9), Inches(3.5), Inches(5.5), Inches(2.5), [
    "Style controls: formal / casual / minimal / sporty",
    "Fit preference: slim, regular, relaxed, oversized",
    "Avoid features: hoods, logos, cropped, skinny fit",
    "Free-text input: \"more formal, avoid skinny fit\"",
], size=17, color=OFFWHITE, spacing=10)

bullets(s, Inches(7.0), Inches(3.5), Inches(5.5), Inches(2.5), [
    "Text prompts encoded with CLIP (e.g. \"a formal outfit\")",
    "Each result gets a goal bonus and avoid penalty",
    "Results re-sorted by final score",
    "Both baseline and reranked results shown side by side",
], size=17, color=OFFWHITE, spacing=10)

key_line(s, "\"Users can now steer results — instead of just accepting them.\"",
         Inches(6.35), ACCENT2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — STAGE 3: FASHIONCLIP + FAISS  (hero slide)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT2)
stage_badge(s, "STAGE 3", ACCENT2)

txt(s, Inches(0.9), Inches(0.55), Inches(10), Inches(0.7),
    "Major Upgrade: FashionCLIP + FAISS", size=36, bold=True, color=WHITE)
divider(s, Inches(1.35), ACCENT2)

txt(s, Inches(0.9), Inches(1.6), Inches(11), Inches(0.55),
    "\"I upgraded both the model and the search system.\"",
    size=20, color=MUTED, italic=True)

# LEFT panel: Model
rect(s, Inches(0.9), Inches(2.4), Inches(5.4), Inches(3.4), PANEL)
rect(s, Inches(0.9), Inches(2.4), Inches(5.4), Inches(0.38), ACCENT)
txt(s, Inches(1.0), Inches(2.44), Inches(5.2), Inches(0.35),
    "MODEL", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, Inches(1.0), Inches(2.95), Inches(2.2), Inches(0.4),
    "Before", size=14, color=MUTED, bold=True)
txt(s, Inches(3.4), Inches(2.95), Inches(2.6), Inches(0.4),
    "After", size=14, color=ACCENT2, bold=True)

txt(s, Inches(1.0), Inches(3.4), Inches(2.0), Inches(0.6),
    "OpenAI CLIP", size=16, color=MUTED)
txt(s, Inches(3.0), Inches(3.4), Inches(0.5), Inches(0.6),
    "\u2192", size=20, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, Inches(3.5), Inches(3.4), Inches(2.5), Inches(0.6),
    "FashionCLIP", size=16, color=ACCENT2, bold=True)

bullets(s, Inches(1.0), Inches(4.1), Inches(5.0), Inches(1.5), [
    "General-purpose → fashion-specific",
    "Trained on 800K fashion image-text pairs",
    "Better at garment types, styles, fabrics",
], size=15, color=OFFWHITE, spacing=8)

# RIGHT panel: Search
rect(s, Inches(7.0), Inches(2.4), Inches(5.4), Inches(3.4), PANEL)
rect(s, Inches(7.0), Inches(2.4), Inches(5.4), Inches(0.38), ACCENT)
txt(s, Inches(7.1), Inches(2.44), Inches(5.2), Inches(0.35),
    "SEARCH", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, Inches(7.1), Inches(2.95), Inches(2.2), Inches(0.4),
    "Before", size=14, color=MUTED, bold=True)
txt(s, Inches(9.5), Inches(2.95), Inches(2.5), Inches(0.4),
    "After", size=14, color=ACCENT2, bold=True)

txt(s, Inches(7.1), Inches(3.4), Inches(2.0), Inches(0.6),
    "NumPy dot product", size=16, color=MUTED)
txt(s, Inches(9.1), Inches(3.4), Inches(0.5), Inches(0.6),
    "→", size=20, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, Inches(9.6), Inches(3.4), Inches(2.5), Inches(0.6),
    "FAISS Index", size=16, color=ACCENT2, bold=True)

bullets(s, Inches(7.1), Inches(4.1), Inches(5.0), Inches(1.5), [
    "Brute-force → production vector search",
    "Index saved to disk, loaded on reuse",
    "Scales to any catalog size",
], size=15, color=OFFWHITE, spacing=8)

key_line(s, "\"This moves the system from a prototype toward a real-world pipeline.\"",
         Inches(6.35), GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT)

txt(s, Inches(0.9), Inches(0.55), Inches(10), Inches(0.7),
    "How It Works", size=36, bold=True)
divider(s, Inches(1.35), ACCENT)

# Pipeline steps
steps = [
    ("Upload\nImage", "Streamlit UI", ACCENT),
    ("FashionCLIP\nEncoder", "512-dim embedding", RGBColor(0x3A, 0x3A, 0x7A)),
    ("FAISS\nSearch", "Top-5 neighbors", RGBColor(0x1A, 0x5A, 0x4A)),
    ("Preference\nReranking", "Goal / Avoid scoring", RGBColor(0x5A, 0x4A, 0x1A)),
    ("Results\nDisplay", "Baseline + Reranked", ACCENT2),
]

bw, bh = Inches(1.9), Inches(1.5)
by = Inches(2.7)
start = Inches(0.65)
gap = Inches(0.52)

for i, (title, sub, color) in enumerate(steps):
    bx = start + i * (bw + gap)
    rect(s, bx, by, bw, bh, color)
    txt(s, bx + Inches(0.1), by + Inches(0.15), bw - Inches(0.2), Inches(0.7),
        title, size=17, bold=True, align=PP_ALIGN.CENTER)
    txt(s, bx + Inches(0.05), by + Inches(0.95), bw - Inches(0.1), Inches(0.45),
        sub, size=12, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = bx + bw + Inches(0.1)
        txt(s, ax, by + Inches(0.45), Inches(0.32), Inches(0.6),
            "\u2192", size=26, color=MUTED, align=PP_ALIGN.CENTER)

# One-liner description per step
descs = [
    "User uploads a\nclothing photo",
    "Image converted to\nnumerical embedding",
    "Nearest catalog\nitems retrieved",
    "Style preferences\napplied to scores",
    "Both views shown\nside by side",
]
for i, desc in enumerate(descs):
    bx = start + i * (bw + gap)
    txt(s, bx, by + bh + Inches(0.2), bw, Inches(0.7),
        desc, size=13, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WHAT I FIXED (repo cleanup)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5), GOLD)

txt(s, Inches(0.9), Inches(0.55), Inches(10), Inches(0.7),
    "Engineering Discipline: What I Fixed", size=36, bold=True)
divider(s, Inches(1.35), GOLD)

txt(s, Inches(0.9), Inches(1.6), Inches(11), Inches(0.55),
    "Not just new features — I also cut what wasn't working.",
    size=20, color=MUTED, italic=True)

# Before column
rect(s, Inches(0.9), Inches(2.4), Inches(5.4), Inches(3.3), PANEL)
txt(s, Inches(1.0), Inches(2.55), Inches(5.0), Inches(0.45),
    "Before", size=17, color=WARN, bold=True)
bullets(s, Inches(1.0), Inches(3.1), Inches(5.0), Inches(2.4), [
    "Broken training pipeline (WardrobeNet missing)",
    "Stubbed evaluate.py with dummy predictions",
    "Unused config files referencing wrong model",
    "Two disconnected systems in one repo",
], size=15, color=OFFWHITE, spacing=10)

# After column
rect(s, Inches(7.0), Inches(2.4), Inches(5.4), Inches(3.3), PANEL)
txt(s, Inches(7.1), Inches(2.55), Inches(5.0), Inches(0.45),
    "After", size=17, color=ACCENT2, bold=True)
bullets(s, Inches(7.1), Inches(3.1), Inches(5.0), Inches(2.4), [
    "Moved broken code to experimental/ (preserved, not deleted)",
    "src/ has only 4 files — all connected to the demo",
    "README matches syllabus requirements exactly",
    "One working pipeline, clearly documented",
], size=15, color=OFFWHITE, spacing=10)

key_line(s, "\"I prioritized a clean, working pipeline over incomplete features.\"",
         Inches(6.35), GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WHAT'S NEXT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT2)

txt(s, Inches(0.9), Inches(0.55), Inches(10), Inches(0.7),
    "What's Next", size=36, bold=True)
divider(s, Inches(1.35), ACCENT2)

txt(s, Inches(0.9), Inches(1.6), Inches(11), Inches(0.55),
    "Realistic next steps — not promises.",
    size=20, color=MUTED, italic=True)

items = [
    ("Persistent Feedback",
     "Store user likes and dislikes across sessions to personalize retrieval over time."),
    ("Larger Catalog",
     "More catalog images = better coverage. Results improve directly with catalog variety."),
    ("Cloud Deployment",
     "Deploy to Streamlit Cloud or Render so the demo runs without a local setup."),
]

for i, (title, desc) in enumerate(items):
    y = Inches(2.6) + Inches(i * 1.4)
    rect(s, Inches(0.9), y, Inches(0.08), Inches(0.9), ACCENT2)
    txt(s, Inches(1.2), y, Inches(4), Inches(0.5),
        title, size=20, bold=True, color=WHITE)
    txt(s, Inches(1.2), y + Inches(0.5), Inches(11), Inches(0.55),
        desc, size=17, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LIVE DEMO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), BG)

txt(s, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
    "Live Demo", size=52, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(3.5), Inches(3.3), Inches(6.3), Inches(0.75), PANEL)
txt(s, Inches(3.6), Inches(3.38), Inches(6.1), Inches(0.6),
    "streamlit run app.py", size=24, color=ACCENT2,
    align=PP_ALIGN.CENTER, font="Consolas", bold=True)

txt(s, Inches(1), Inches(4.4), Inches(11), Inches(0.6),
    "Upload  →  Retrieve  →  Preference Filter  →  Reranked Results",
    size=18, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, Inches(0), Inches(0), Inches(0.06), Inches(7.5), ACCENT)

txt(s, Inches(1), Inches(2.5), Inches(11), Inches(1.0),
    "Thank You", size=52, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(4.0), Inches(3.8), Inches(5.3), Inches(0.04), ACCENT)

txt(s, Inches(1), Inches(4.1), Inches(11), Inches(0.5),
    "Questions?", size=26, color=MUTED, align=PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
    "github.com/yeobian/Capstone",
    size=16, color=MUTED, align=PP_ALIGN.CENTER, font="Consolas")


prs.save("Capstone_Presentation.pptx")
print("Saved: Capstone_Presentation.pptx  —  10 slides")
